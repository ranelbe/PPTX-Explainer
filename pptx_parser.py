from pptx import Presentation, exc

# constants error messages
FILE_ERROR = "File doesn't exist, or file is not in the correct format."
SLIDE_PARSE_ERROR = "Cannot parse this slide: "


class PPTXParser:
    """ Parse the presentation to get its data """

    def __init__(self, presentation_path):
        """
           initialization.
           :param presentation_path: path to the PowerPoint presentation.
        """
        self._presentation_path = presentation_path
        self._presentation = None
        self._slides = []

    def extract_slides_text(self):
        """
        Extracts text from a PowerPoint presentation.
        :return: list of slides text.
        """
        self.open_presentation()
        self._slides = [self._extract_slide_text(slide)
                        for slide in self._presentation.slides]
        return self._slides

    def open_presentation(self):
        """
        Opens a PowerPoint presentation.
        :raises: ValueError if the presentation doesn't exist.
        or is not in the correct format.
        """
        try:
            self._presentation = Presentation(self._presentation_path)
        except exc.PackageNotFoundError as e:
            raise ValueError(f"{FILE_ERROR}\n{e}")

    @staticmethod
    def _extract_slide_text(slide):
        """
        Extracts text from the given slide.
        :param slide: slide to extract text from.
        :return: the slide text.
        :raises: ValueError if the slide cannot be parsed.
        """
        text_runs = []
        try:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_runs.append(run.text.strip())
                        # add a newline after each paragraph
                        if paragraph.runs:
                            text_runs.append('\n')
        except Exception as e:
            return f"{SLIDE_PARSE_ERROR} {e}"

        return ' '.join(text_runs)
